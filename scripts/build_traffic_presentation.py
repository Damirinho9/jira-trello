#!/usr/bin/env python3
"""Build task #4 PowerPoint slide using python-pptx + direct XML animation editing.

The script draws the scene from the methodical guide (vertical arc road, car, traffic
light) and injects animation timing XML mapped to created shape IDs.
"""

from __future__ import annotations

import argparse
import pathlib
import zipfile
from copy import deepcopy
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu, Pt

P_NAMESPACE = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS = {"p": P_NAMESPACE}

# Shape IDs used in bundled template animation timeline (task4_traffic.pptx).
TEMPLATE_SHAPE_IDS = {
    "car": "7",
    "red": "10",
    "yellow": "11",
    "green": "12",
}

# Motion paths updated to match the vertical trajectory from the textbook.
MOTION_PATHS = (
    "M 0 0 C -0.03 0.18 -0.08 0.42 -0.13 0.58",
    "M 0 0 C -0.04 0.14 -0.09 0.28 -0.13 0.42",
)


class PresentationBuilder:
    def __init__(self, output_path: pathlib.Path, timing_template: pathlib.Path) -> None:
        self.output_path = output_path
        self.timing_template = timing_template

    def build(self) -> None:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        scene = self._draw_scene(slide, prs.slide_width, prs.slide_height)
        self._inject_timing_xml(slide, scene)

        prs.save(self.output_path)

    def _draw_scene(self, slide, slide_w: Emu, slide_h: Emu) -> dict[str, object]:
        # Vertical arc road (as in методичка: path from top to bottom, curved left).
        road = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ARC,
            int(slide_w * 0.12),
            int(slide_h * 0.05),
            int(slide_w * 0.58),
            int(slide_h * 0.9),
        )
        road.fill.background()
        road.line.color.rgb = RGBColor(95, 95, 95)
        road.line.width = Pt(38)
        road.rotation = 92

        car = self._draw_car(slide, slide_w, slide_h)
        traffic = self._draw_traffic_light(slide, slide_w, slide_h)

        return {
            "road": road,
            "car": car,
            "red": traffic["red"],
            "yellow": traffic["yellow"],
            "green": traffic["green"],
        }

    def _draw_car(self, slide, slide_w: Emu, slide_h: Emu):
        body = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(slide_w * 0.16),
            int(slide_h * 0.10),
            int(slide_w * 0.16),
            int(slide_h * 0.065),
        )
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(35, 35, 35)
        body.line.fill.background()

        roof = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(slide_w * 0.19),
            int(slide_h * 0.075),
            int(slide_w * 0.095),
            int(slide_h * 0.04),
        )
        roof.fill.solid()
        roof.fill.fore_color.rgb = RGBColor(55, 55, 55)
        roof.line.fill.background()

        wheel_l = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            int(slide_w * 0.18),
            int(slide_h * 0.145),
            int(slide_w * 0.035),
            int(slide_w * 0.035),
        )
        wheel_r = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            int(slide_w * 0.255),
            int(slide_h * 0.145),
            int(slide_w * 0.035),
            int(slide_w * 0.035),
        )
        for wheel in (wheel_l, wheel_r):
            wheel.fill.solid()
            wheel.fill.fore_color.rgb = RGBColor(10, 10, 10)
            wheel.line.fill.background()

        car = slide.shapes.add_group_shape([body, roof, wheel_l, wheel_r])
        car.rotation = -18
        return car

    def _draw_traffic_light(self, slide, slide_w: Emu, slide_h: Emu) -> dict[str, object]:
        pole = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            int(slide_w * 0.67),
            int(slide_h * 0.28),
            int(slide_w * 0.012),
            int(slide_h * 0.44),
        )
        pole.fill.solid()
        pole.fill.fore_color.rgb = RGBColor(35, 35, 35)
        pole.line.fill.background()

        housing = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(slide_w * 0.63),
            int(slide_h * 0.22),
            int(slide_w * 0.09),
            int(slide_h * 0.24),
        )
        housing.fill.solid()
        housing.fill.fore_color.rgb = RGBColor(60, 60, 60)
        housing.line.fill.background()

        red = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            int(slide_w * 0.648),
            int(slide_h * 0.245),
            int(slide_w * 0.048),
            int(slide_w * 0.048),
        )
        yellow = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            int(slide_w * 0.648),
            int(slide_h * 0.31),
            int(slide_w * 0.048),
            int(slide_w * 0.048),
        )
        green = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            int(slide_w * 0.648),
            int(slide_h * 0.375),
            int(slide_w * 0.048),
            int(slide_w * 0.048),
        )

        for lamp in (red, yellow, green):
            lamp.fill.solid()
            lamp.fill.fore_color.rgb = RGBColor(85, 85, 85)
            lamp.line.fill.background()

        return {"pole": pole, "housing": housing, "red": red, "yellow": yellow, "green": green}

    def _inject_timing_xml(self, slide, scene: dict[str, object]) -> None:
        timing = self._load_template_timing()

        new_ids = {
            "car": str(scene["car"].shape_id),
            "red": str(scene["red"].shape_id),
            "yellow": str(scene["yellow"].shape_id),
            "green": str(scene["green"].shape_id),
        }

        for sp_tgt in timing.findall(".//p:spTgt", NS):
            old = sp_tgt.attrib.get("spid")
            for key, old_id in TEMPLATE_SHAPE_IDS.items():
                if old == old_id:
                    sp_tgt.set("spid", new_ids[key])

        motion_nodes = timing.findall(".//p:animMotion", NS)
        if len(motion_nodes) >= 2:
            motion_nodes[0].set("path", MOTION_PATHS[0])
            motion_nodes[1].set("path", MOTION_PATHS[1])

        slide._element.append(timing)

    def _load_template_timing(self):
        with zipfile.ZipFile(self.timing_template) as zf:
            xml = zf.read("ppt/slides/slide1.xml")

        root = ET.fromstring(xml)
        timing = root.find("p:timing", NS)
        if timing is None:
            raise RuntimeError(f"No <p:timing> section found in template: {self.timing_template}")

        return deepcopy(timing)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Build task #4 traffic animation PPTX using python-pptx + XML timing"
    )
    parser.add_argument(
        "--output",
        type=pathlib.Path,
        default=pathlib.Path("task4_traffic.pptx"),
        help="Output pptx file path",
    )
    parser.add_argument(
        "--timing-template",
        type=pathlib.Path,
        default=pathlib.Path("task4_traffic.pptx"),
        help="PPTX file used only as timing XML template",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    builder = PresentationBuilder(args.output, args.timing_template)
    builder.build()
    print(f"Presentation saved to: {args.output}")


if __name__ == "__main__":
    main()
